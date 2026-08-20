import os
import re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = "/Users/metrobee/Library/CloudStorage/GoogleDrive-borismeldre@gmail.com/My Drive/myGslides/fungib_hundiallika.pptx"
INDEX_HTML_PATH = "/Users/metrobee/Projects/fungib/public/index.html"
PUBLIC_DIR = "/Users/metrobee/Projects/fungib/public"

def parse_pptx_slides():
    if not os.path.exists(PPTX_PATH):
        print(f"[VIGA] PowerPoint faili ei leitud: {PPTX_PATH}")
        return None
        
    prs = Presentation(PPTX_PATH)
    slides_data = []
    
    for idx, slide in enumerate(prs.slides):
        slide_info = {
            'title': '',
            'subtitle': '',
            'meta': '',
            'bullets': [],
            'image': None,
            'is_title': idx == 0
        }
        
        text_shapes = []
        image_shape = None
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_shapes.append(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_shape = shape
                
        # Sort text shapes by top position
        text_shapes.sort(key=lambda s: s.top)
        
        if text_shapes:
            slide_info['title'] = text_shapes[0].text.strip()
            
            if len(text_shapes) > 1:
                for body_shape in text_shapes[1:]:
                    lines = [line.strip() for line in body_shape.text.split('\n') if line.strip()]
                    
                    if slide_info['is_title']:
                        if lines:
                            if not slide_info['subtitle']:
                                slide_info['subtitle'] = lines[0]
                                extra_lines = lines[1:]
                            else:
                                extra_lines = lines
                            if extra_lines:
                                if slide_info['meta']:
                                    slide_info['meta'] += " | " + " | ".join(extra_lines)
                                else:
                                    slide_info['meta'] = " | ".join(extra_lines)
                    else:
                        for line in lines:
                            clean_line = re.sub(r'^[•\-\*\s]+', '', line).strip()
                            if clean_line:
                                if clean_line.startswith('"') and clean_line.endswith('"'):
                                    slide_info['bullets'].append(clean_line)
                                elif len(slide_info['bullets']) == 0 and not line.startswith('•') and not slide_info['subtitle']:
                                    slide_info['subtitle'] = clean_line
                                else:
                                    slide_info['bullets'].append(clean_line)
                                
        # Extract Image if present
        if image_shape:
            try:
                image = image_shape.image
                image_bytes = image.blob
                ext = image.ext
                
                img_name = f"extracted_image_{idx}.{ext}"
                img_path = os.path.join(PUBLIC_DIR, img_name)
                
                with open(img_path, 'wb') as f_img:
                    f_img.write(image_bytes)
                    
                slide_info['image'] = img_name
                print(f"Eraldasin slaidilt {idx+1} pildi ja salvestasin: {img_name}")
            except Exception as e:
                print(f"Hoiatus: Pildi eraldamisel slaidilt {idx+1} tekkis viga: {e}")
                
        slides_data.append(slide_info)
        
    return slides_data

def generate_slide_html(s, idx):
    html = ""
    classes = ["slide"]
    if idx == 0:
        classes.append("active")
    if s['is_title']:
        classes.append("slide-title-card")
        
    class_str = " ".join(classes)
    html += f'                <!-- Slaid {idx+1}: {s["title"]} -->\n'
    html += f'                <div class="{class_str}">\n'
    
    if s['is_title']:
        html += f'                    <h1>{s["title"]}</h1>\n'
        if s['subtitle']:
            html += f'                    <h2>{s["subtitle"]}</h2>\n'
        if s['meta']:
            html += f'                    <div class="lecture-meta">\n'
            parts = [p.strip() for p in s['meta'].split('|')]
            for part in parts:
                if "Lektor" in part:
                    html += f'                        <p style="font-weight: 500; font-size: 1.2rem; color: var(--accent-color);">{part}</p>\n'
                else:
                    html += f'                        <p style="font-size: 0.95rem; color: #64748b; margin-top: 5px;">{part}</p>\n'
            html += f'                    </div>\n'
    else:
        html += f'                    <h2>{s["title"]}</h2>\n'
        if s['subtitle']:
            html += f'                    <p>{s["subtitle"]}</p>\n'
            
        if s['image']:
            html += f'                    <div class="slide-split">\n'
            html += f'                        <div class="slide-split-left">\n'
            
        bullets_html = ""
        has_blockquote = False
        bq_text = ""
        list_items = []
        
        for bullet in s['bullets']:
            if bullet.startswith('"') and bullet.endswith('"'):
                has_blockquote = True
                bq_text = bullet
            else:
                list_items.append(bullet)
                
        if has_blockquote:
            bullets_html += f'                    <blockquote>\n'
            bullets_html += f'                        {bq_text}\n'
            bullets_html += f'                    </blockquote>\n'
            
        if list_items:
            bullets_html += f'                    <ul>\n'
            for item in list_items:
                if ":" in item:
                    parts = item.split(":", 1)
                    bullets_html += f'                        <li><span class="highlight">{parts[0]}:</span>{parts[1]}</li>\n'
                else:
                    bullets_html += f'                        <li>{item}</li>\n'
            bullets_html += f'                    </ul>\n'
            
        if s['image']:
            bullets_html += f'                        </div>\n'
            bullets_html += f'                        <div class="slide-split-right">\n'
            bullets_html += f'                            <img class="slide-img-preview" src="{s["image"]}" alt="{s["title"]}">\n'
            bullets_html += f'                        </div>\n'
            bullets_html += f'                    </div>\n'
            
        html += bullets_html
        
    html += f'                </div>\n\n'
    return html

def update_html():
    slides_data = parse_pptx_slides()
    if not slides_data:
        return
        
    print(f"Genereerin uut HTML-i {len(slides_data)} slaidi jaoks...")
    
    # Build slides HTML string
    slides_html = ""
    for idx, s in enumerate(slides_data):
        slides_html += generate_slide_html(s, idx)
        
    # Read index.html
    with open(INDEX_HTML_PATH, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
        
    # Find slide-deck container and replace its content
    deck = soup.find('div', class_='slide-deck')
    if not deck:
        print("[VIGA] Failist index.html ei leitud div elementiklassiga 'slide-deck'.")
        return
        
    # Replace content of slide-deck
    # We parse the slides_html string into a BeautifulSoup object and insert it
    new_deck_soup = BeautifulSoup(slides_html, 'html.parser')
    deck.clear()
    deck.append(new_deck_soup)
    
    # Update total-slides-num span
    total_span = soup.find('span', id='total-slides-num')
    if total_span:
        total_span.string = str(len(slides_data))
        print(f"Uuendasin slaidide koguarvuks indikaatoris: {len(slides_data)}")
        
    # Save index.html back
    with open(INDEX_HTML_PATH, 'w', encoding='utf-8') as f:
        # Save prettified or raw (prettify might format the whole doc, let's save as unicode string)
        f.write(str(soup))
        
    print(f"[OK] HTML fail edukalt sünkroniseeritud PowerPointiga: {INDEX_HTML_PATH}")

if __name__ == "__main__":
    update_html()
