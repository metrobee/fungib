import os
import sys
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

INDEX_HTML_PATH = "/Users/metrobee/Projects/fungib/public/index.html"
PUBLIC_DIR = "/Users/metrobee/Projects/fungib/public"
OUTPUT_DIR = "/Users/metrobee/Library/CloudStorage/GoogleDrive-borismeldre@gmail.com/My Drive/myGslides"
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "fungib_hundiallika.pptx")

# Slidesgo 2026 Deep Teal colors
DARK_GREEN_BG = RGBColor(26, 54, 58) # #1a363a (Deep Teal background)
LIGHT_GREEN_ACCENT = RGBColor(46, 196, 182) # #2ec4b6 (Transformative Teal accent)
WHITE_COLOR = RGBColor(255, 255, 255)
SUB_TEXT_COLOR = RGBColor(214, 216, 210) # #d6d8d2 (Wax Paper cream)

def parse_html_slides():
    with open(INDEX_HTML_PATH, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
        
    slides_data = []
    slide_divs = soup.find_all('div', class_='slide')
    
    for div in slide_divs:
        slide = {
            'is_title': 'slide-title-card' in div.get('class', []),
            'title': '',
            'subtitle': '',
            'meta': '',
            'bullets': [],
            'image': None
        }
        
        # Extract title
        h1 = div.find('h1')
        h2 = div.find('h2')
        if h1:
            slide['title'] = h1.get_text().strip()
        elif h2:
            slide['title'] = h2.get_text().strip()
            
        # Extract subtitle / meta / description
        if slide['is_title']:
            h2_sub = div.find('h2')
            if h2_sub and h1:
                slide['subtitle'] = h2_sub.get_text().strip()
            meta_p = div.find('div', class_='lecture-meta')
            if meta_p:
                slide['meta'] = meta_p.get_text(" | ").strip()
        else:
            p_desc = div.find('p')
            if p_desc:
                slide['subtitle'] = p_desc.get_text().strip()
                
            # Check for blockquote
            bq = div.find('blockquote')
            if bq:
                slide['bullets'].append(f'"{bq.get_text().strip()}"')
                
            # Check for list items
            lis = div.find_all('li')
            for li in lis:
                slide['bullets'].append(li.get_text().strip())
                
            # Check for image
            img = div.find('img')
            if img:
                src = img.get('src')
                slide['image'] = src
                
        slides_data.append(slide)
    return slides_data

def build_pptx():
    print("Parsing slides from index.html...")
    slides_data = parse_html_slides()
    print(f"Parsed {len(slides_data)} slides.")
    
def build_pptx():
    print("Parsing slides from index.html...")
    slides_data = parse_html_slides()
    print(f"Parsed {len(slides_data)} slides.")
    
    # Ensure output directory exists
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    TEMPLATE_PATH = "/Users/metrobee/Projects/fungib/template.pptx"
    using_template = os.path.exists(TEMPLATE_PATH)
    
    if using_template:
        print(f"Found ready-made template: {TEMPLATE_PATH}. Building presentation on top of it...")
        prs = Presentation(TEMPLATE_PATH)
        # Clear all placeholder slides in the template
        slide_ids = prs.slides._sldIdLst
        while len(slide_ids) > 0:
            prs.part.drop_rel(slide_ids[0].rId)
            del slide_ids[0]
            
        print(f"Template slide layouts available: {len(prs.slide_layouts)}")
        
        # Detect if it's a light-themed template (e.g. minimalist, clean) or dark-themed
        is_light_template = "minimalist" in os.path.basename(TEMPLATE_PATH).lower() or "light" in os.path.basename(TEMPLATE_PATH).lower()
        
        if is_light_template:
            # Light theme colors for minimalist template
            t_bg = None
            t_title_color = RGBColor(31, 39, 34)       # Charcoal
            t_sub_color = RGBColor(96, 106, 100)       # Olive Muted Gray
            t_body_color = RGBColor(78, 86, 81)        # Dark grey
            t_accent_color = RGBColor(176, 125, 80)    # Terracotta Accent
            t_meta_color = RGBColor(128, 136, 131)     # Muted grey
        else:
            # Dark theme colors (like the Fungi Lesson template)
            t_bg = None
            t_title_color = RGBColor(255, 255, 255)    # White title (or keep it white for contrast)
            t_sub_color = RGBColor(229, 193, 88)       # Gold/sand accent (#e5c158)
            t_body_color = RGBColor(255, 255, 255)     # White body
            t_accent_color = RGBColor(46, 196, 182)    # Teal accent
            t_meta_color = RGBColor(214, 216, 210)     # Wax paper cream
    else:
        print("No template.pptx found. Building custom presentation from scratch...")
        prs = Presentation()
        # Set to widescreen 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        t_bg = DARK_GREEN_BG
        t_title_color = LIGHT_GREEN_ACCENT
        t_sub_color = WHITE_COLOR
        t_body_color = WHITE_COLOR
        t_accent_color = LIGHT_GREEN_ACCENT
        t_meta_color = SUB_TEXT_COLOR
        
    for idx, s in enumerate(slides_data):
        # Choose layout from template if available
        if using_template:
            # Layout 0 is usually Title layout. Layout 6 is usually Blank.
            # We try using Layout 0 for title, and Layout 6 (or whatever is blank/content) for rest
            if s['is_title']:
                layout = prs.slide_layouts[0]
            else:
                # Use layout 6 (Blank) if available, otherwise layout 1 or 0
                layout_idx = 6 if len(prs.slide_layouts) > 6 else (1 if len(prs.slide_layouts) > 1 else 0)
                layout = prs.slide_layouts[layout_idx]
        else:
            layout = prs.slide_layouts[6] # Blank layout in new prs
            
        slide = prs.slides.add_slide(layout)
        
        # Remove default placeholders so they don't overlay our content with "Click to add title"
        if using_template:
            for shape in list(slide.shapes):
                if shape.is_placeholder:
                    try:
                        slide.shapes.element.remove(shape.element)
                    except Exception as e:
                        print(f"Could not remove placeholder: {e}")
        
        # Set slide background color only if we are NOT using a template
        if not using_template:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = t_bg
            
        # Helper to apply font properties cleanly and inherit theme font
        def style_paragraph(p, size, color, bold=False, italic=False, space_after=0):
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.italic = italic
            if space_after:
                p.space_after = Pt(space_after)
            if not using_template:
                p.font.name = 'Arial'

        # Add Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = Inches(0)
        tf_title.margin_right = Inches(0)
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = s['title']
        style_paragraph(p_title, size=36 if s['is_title'] else 28, color=t_title_color, bold=True)
        
        # Add Content Box
        if s['is_title']:
            # Title Slide metadata/subtitle
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = Inches(0)
            tf_content.margin_right = Inches(0)
            
            p_sub = tf_content.paragraphs[0]
            p_sub.text = s['subtitle']
            style_paragraph(p_sub, size=20, color=t_sub_color, space_after=24)
            
            if s['meta']:
                p_meta = tf_content.add_paragraph()
                p_meta.text = s['meta']
                style_paragraph(p_meta, size=14, color=t_meta_color)
        else:
            # Regular slide content
            content_w = Inches(6.5) if s['image'] else Inches(11.7)
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), content_w, Inches(4.8))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = Inches(0)
            tf_content.margin_right = Inches(0)
            
            first_p = True
            if s['subtitle']:
                p_sub = tf_content.paragraphs[0]
                p_sub.text = s['subtitle']
                style_paragraph(p_sub, size=16, color=t_sub_color, space_after=16)
                first_p = False
                
            for bullet in s['bullets']:
                if first_p:
                    p = tf_content.paragraphs[0]
                    first_p = False
                else:
                    p = tf_content.add_paragraph()
                
                # Check if blockquote (quote marks around it)
                if bullet.startswith('"') and bullet.endswith('"'):
                    p.text = bullet
                    style_paragraph(p, size=18, color=t_accent_color, italic=True, space_after=16)
                else:
                    p.text = "• " + bullet
                    style_paragraph(p, size=15, color=t_body_color, space_after=12)
                
            # Add Image if present
            if s['image']:
                img_path = os.path.join(PUBLIC_DIR, s['image'])
                if os.path.exists(img_path):
                    img_left = Inches(7.8)
                    img_top = Inches(1.8)
                    img_width = Inches(4.7)
                    img_height = Inches(4.5)
                    slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
                    print(f"Added image {s['image']} to slide.")
                else:
                    print(f"Image path not found: {img_path}")
                    
    prs.save(OUTPUT_PATH)
    print(f"\n[OK] PowerPoint esitlus salvestatud: {OUTPUT_PATH}")

if __name__ == "__main__":
    build_pptx()
